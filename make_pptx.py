import pptx
import glob
p = pptx.Presentation()
ims = sorted(glob.glob("ims/*"))
blank_slide_layout= p.slide_layouts[6]
for im in ims:
    slide = p.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(im, 0, 0, p.slide_width, p.slide_height)

p.save("test.pptx")
